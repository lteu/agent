"""
技术汇报 PPT 模板生成器
主题：AI Agent 技术汇报
风格：深色科技感 + 渐变强调色
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ===== 配色方案 =====
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 深蓝黑背景
BG_CARD = RGBColor(0x1A, 0x25, 0x40)        # 卡片背景
ACCENT_PRIMARY = RGBColor(0x00, 0xD4, 0xFF)  # 主强调色 - 青蓝
ACCENT_SECOND = RGBColor(0x7B, 0x61, 0xFF)   # 次强调色 - 紫色
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # 主文字
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)       # 次要文字
TEXT_LIGHT = RGBColor(0xC9, 0xD1, 0xD9)      # 三级文字
DIVIDER = RGBColor(0x2D, 0x3A, 0x5C)         # 分割线
HIGHLIGHT = RGBColor(0xFF, 0xD1, 0x66)       # 高亮色 - 金色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank_layout = prs.slide_layouts[6]


def add_gradient_rect(slide, left, top, width, height, color1, color2, angle=45):
    """添加渐变填充矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    # 设置渐变角度
    fill.gradient_angle = angle
    return shape


def set_bg(slide, color=BG_DARK):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=TEXT_WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # 设置中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    return txBox


def add_rich_text(slide, left, top, width, height, paragraphs_data, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    """添加富文本（多段落多格式）
    paragraphs_data: list of list of dict，每个段落由多个run组成
    run dict keys: text, size, bold, color, font
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    for pi, para_runs in enumerate(paragraphs_data):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)

        for run_data in para_runs:
            run = p.add_run()
            run.text = run_data.get("text", "")
            run.font.size = Pt(run_data.get("size", 16))
            run.font.bold = run_data.get("bold", False)
            run.font.color.rgb = run_data.get("color", TEXT_WHITE)
            run.font.name = run_data.get("font", font)
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', run_data.get("font", font))

    return txBox


def add_card(slide, left, top, width, height, radius=0.1, fill=BG_CARD, border_color=None):
    """添加圆角卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_decorative_line(slide, left, top, width, color=ACCENT_PRIMARY, thickness=3):
    """添加装饰线条"""
    line = slide.shapes.add_connector(1, left, top, left + width, top)  # 1 = straight
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_page_number(slide, num, total):
    """添加页码"""
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num} / {total}", size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)


def add_footer_brand(slide):
    """添加底部品牌标识"""
    add_text(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.4),
             "AI Agent Technical Report", size=10, color=TEXT_GRAY)


# ==========================================
# 第1页：封面
# ==========================================
def make_cover():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # 左侧装饰渐变条
    add_gradient_rect(slide, Inches(0), Inches(0), Inches(0.15), SH,
                      ACCENT_PRIMARY, ACCENT_SECOND, angle=90)

    # 右上角装饰圆（光晕效果）
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5),
                                   Inches(4), Inches(4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_PRIMARY
    glow.fill.fore_color.brightness = 0.6
    glow.line.fill.background()
    # 透明度
    sp = glow._element
    srgb = sp.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '15000')

    # 底部装饰线条组
    for i, (w, c) in enumerate([(120, ACCENT_PRIMARY), (80, ACCENT_SECOND), (50, HIGHLIGHT)]):
        add_decorative_line(slide, Inches(0.8 + i * 1.5), Inches(6.2),
                           Inches(0.8 + i * 0.3), color=c, thickness=4)

    # 标签
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), Inches(1.8), Inches(2.2), Inches(0.5))
    tag.adjustments[0] = 0.5
    tag.fill.solid()
    tag.fill.fore_color.rgb = BG_CARD
    tag.line.color.rgb = ACCENT_PRIMARY
    tag.line.width = Pt(1)
    tf = tag.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TECHNICAL REPORT"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_PRIMARY
    run.font.name = "Arial"

    # 主标题
    add_text(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
             "AI Agent 技术架构与实践", size=52, bold=True, color=TEXT_WHITE)

    # 副标题
    add_text(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.8),
             "从原理到落地：智能体的核心能力与工程实践", size=22, color=TEXT_GRAY)

    # 分割线
    add_decorative_line(slide, Inches(0.8), Inches(4.9), Inches(2.5),
                       color=ACCENT_PRIMARY, thickness=3)

    # 信息行
    add_rich_text(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(1),
                  [
                      [{"text": "汇报人：", "size": 14, "color": TEXT_GRAY},
                       {"text": "技术团队", "size": 14, "color": TEXT_LIGHT}],
                      [{"text": "日  期：", "size": 14, "color": TEXT_GRAY},
                       {"text": "2025 年 7 月", "size": 14, "color": TEXT_LIGHT}],
                  ])

    add_page_number(slide, 1, 12)
    return slide


# ==========================================
# 第2页：目录
# ==========================================
def make_toc():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # 左侧渐变条
    add_gradient_rect(slide, Inches(0), Inches(0), Inches(0.15), SH,
                      ACCENT_PRIMARY, ACCENT_SECOND, angle=90)

    # 标题区
    add_text(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.6),
             "CONTENTS", size=14, bold=True, color=ACCENT_PRIMARY, font="Arial")
    add_text(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(1),
             "目  录", size=36, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.8), Inches(1.95), Inches(1.5),
                       color=ACCENT_PRIMARY, thickness=3)

    # 目录项 - 两列布局
    items = [
        ("01", "背景与趋势", "AI Agent 发展背景与行业趋势"),
        ("02", "核心概念", "Agent 的定义、组成与核心能力"),
        ("03", "技术架构", "整体架构设计与关键模块解析"),
        ("04", "关键技术", "规划、记忆、工具调用等核心技术"),
        ("05", "实践案例", "典型落地场景与效果数据"),
        ("06", "挑战与展望", "当前瓶颈与未来发展方向"),
    ]

    col_x = [Inches(0.8), Inches(6.8)]
    for i, (num, title, desc) in enumerate(items):
        col = i // 3
        row = i % 3
        x = col_x[col]
        y = Inches(2.6 + row * 1.5)

        # 卡片
        card = add_card(slide, x, y, Inches(5.8), Inches(1.2))

        # 编号
        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.8),
                 num, size=32, bold=True, color=ACCENT_PRIMARY, font="Arial",
                 anchor=MSO_ANCHOR.MIDDLE)

        # 标题
        add_text(slide, x + Inches(1.2), y + Inches(0.2), Inches(4.3), Inches(0.5),
                 title, size=18, bold=True, color=TEXT_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

        # 描述
        add_text(slide, x + Inches(1.2), y + Inches(0.65), Inches(4.3), Inches(0.45),
                 desc, size=12, color=TEXT_GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer_brand(slide)
    add_page_number(slide, 2, 12)
    return slide


# ==========================================
# 第3页：章节过渡页
# ==========================================
def make_section(num, title, subtitle):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # 左侧大号章节号
    add_text(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(3),
             num, size=200, bold=True, color=BG_CARD, font="Arial",
             anchor=MSO_ANCHOR.MIDDLE)

    # 右侧内容
    # 标签
    add_text(slide, Inches(6.5), Inches(2.5), Inches(5), Inches(0.5),
             "CHAPTER " + num, size=14, bold=True, color=ACCENT_PRIMARY, font="Arial")

    # 装饰线
    add_decorative_line(slide, Inches(6.5), Inches(3.0), Inches(1.2),
                       color=ACCENT_PRIMARY, thickness=3)

    # 章节标题
    add_text(slide, Inches(6.5), Inches(3.3), Inches(6), Inches(1),
             title, size=40, bold=True, color=TEXT_WHITE)

    # 副标题
    add_text(slide, Inches(6.5), Inches(4.3), Inches(6), Inches(0.8),
             subtitle, size=16, color=TEXT_GRAY)

    # 底部装饰
    add_decorative_line(slide, Inches(6.5), Inches(5.5), Inches(5),
                       color=DIVIDER, thickness=1)

    add_page_number(slide, int(num) + 2, 12)
    return slide


# ==========================================
# 第4页：两栏内容页 - 背景与趋势
# ==========================================
def make_two_column():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # 顶部渐变条
    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    # 页面标题
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
             "AI Agent 发展背景与行业趋势", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)

    # 左栏 - 背景
    left_x = Inches(0.6)
    col_w = Inches(6.0)

    card_l = add_card(slide, left_x, Inches(1.6), col_w, Inches(5.3))

    add_text(slide, left_x + Inches(0.4), Inches(1.85), col_w, Inches(0.5),
             "▍发展背景", size=20, bold=True, color=ACCENT_PRIMARY)

    # 要点列表
    points_left = [
        ("大模型能力跃迁", "LLM 从对话走向推理、规划，为 Agent 提供核心大脑"),
        ("产业需求升级", "从信息检索到任务执行，企业需要端到端的自动化解决方案"),
        ("工具生态成熟", "API 经济、插件体系、RPA 等基础设施日趋完善"),
        ("成本快速下降", "推理成本指数级降低，Agent 规模化部署成为可能"),
    ]

    for i, (title, desc) in enumerate(points_left):
        y = Inches(2.6 + i * 1.05)
        # 圆点标记
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      left_x + Inches(0.4), y + Inches(0.15),
                                      Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_PRIMARY
        dot.line.fill.background()

        add_text(slide, left_x + Inches(0.75), y, Inches(5), Inches(0.4),
                 title, size=15, bold=True, color=TEXT_WHITE)
        add_text(slide, left_x + Inches(0.75), y + Inches(0.4), Inches(5), Inches(0.6),
                 desc, size=12, color=TEXT_GRAY)

    # 右栏 - 趋势
    right_x = Inches(6.9)

    card_r = add_card(slide, right_x, Inches(1.6), col_w, Inches(5.3))

    add_text(slide, right_x + Inches(0.4), Inches(1.85), col_w, Inches(0.5),
             "▍行业趋势", size=20, bold=True, color=ACCENT_SECOND)

    points_right = [
        ("从单 Agent 到多 Agent 协作", "多智能体协同解决复杂任务，效率显著提升"),
        ("垂直领域深耕", "金融、医疗、法律等专业领域 Agent 快速发展"),
        ("端云协同架构", "端侧轻量 Agent + 云端大模型，兼顾隐私与能力"),
        ("自主进化能力", "Agent 具备自我反思与持续学习的进化机制"),
    ]

    for i, (title, desc) in enumerate(points_right):
        y = Inches(2.6 + i * 1.05)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      right_x + Inches(0.4), y + Inches(0.15),
                                      Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_SECOND
        dot.line.fill.background()

        add_text(slide, right_x + Inches(0.75), y, Inches(5), Inches(0.4),
                 title, size=15, bold=True, color=TEXT_WHITE)
        add_text(slide, right_x + Inches(0.75), y + Inches(0.4), Inches(5), Inches(0.6),
                 desc, size=12, color=TEXT_GRAY)

    add_footer_brand(slide)
    add_page_number(slide, 4, 12)
    return slide


# ==========================================
# 第5页：架构图页 - 核心架构
# ==========================================
def make_architecture():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "AI Agent 核心架构", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "六大核心模块构成完整的智能体闭环", size=14, color=TEXT_GRAY)

    # 中心：大模型
    center_x = Inches(5.66)
    center_y = Inches(3.4)
    center_w = Inches(2.0)
    center_h = Inches(1.4)

    core = add_gradient_rect(slide, center_x, center_y, center_w, center_h,
                             ACCENT_PRIMARY, ACCENT_SECOND, angle=45)
    core.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = core.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "大语言模型\n(LLM Core)"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BG_DARK
    run.font.name = "Microsoft YaHei"

    # 环绕模块 - 6个
    import math
    modules = [
        ("感知模块", "Perception", ACCENT_PRIMARY),
        ("规划模块", "Planning", ACCENT_PRIMARY),
        ("记忆模块", "Memory", ACCENT_SECOND),
        ("工具调用", "Tool Use", ACCENT_SECOND),
        ("行动执行", "Action", HIGHLIGHT),
        ("反思学习", "Reflection", HIGHLIGHT),
    ]

    radius_x = Inches(4.2)
    radius_y = Inches(1.8)

    for i, (name, en, color) in enumerate(modules):
        angle = math.radians(i * 60 - 90)
        cx = center_x + center_w / 2 + radius_x * math.cos(angle) - Inches(1.1)
        cy = center_y + center_h / 2 + radius_y * math.sin(angle) - Inches(0.55)

        mod = add_card(slide, cx, cy, Inches(2.2), Inches(1.1),
                       fill=BG_CARD, border_color=color)

        tf = mod.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = color
        r1.font.name = "Microsoft YaHei"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = en
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_GRAY
        r2.font.name = "Arial"

    # 底部说明
    bottom_card = add_card(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.8),
                           fill=BG_CARD)
    add_rich_text(slide, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
                  [[
                      {"text": "闭环机制：", "size": 13, "bold": True, "color": ACCENT_PRIMARY},
                      {"text": "感知 → 规划 → 行动 → 反思 → 记忆更新 → 再感知", "size": 13, "color": TEXT_LIGHT},
                      {"text": "，形成持续进化的智能体循环。", "size": 13, "color": TEXT_GRAY},
                  ]], anchor=MSO_ANCHOR.MIDDLE)

    add_footer_brand(slide)
    add_page_number(slide, 5, 12)
    return slide


# ==========================================
# 第6页：三栏卡片页 - 关键技术
# ==========================================
def make_three_cards():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "三大核心关键技术", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "规划能力、记忆系统、工具使用构成 Agent 能力铁三角", size=14, color=TEXT_GRAY)

    cards = [
        {
            "icon": "🧠",
            "title": "规划能力",
            "subtitle": "Planning & Reasoning",
            "color": ACCENT_PRIMARY,
            "points": [
                "任务拆解与子目标生成",
                "多步推理与路径规划",
                "动态调整与容错机制",
                "ReAct / CoT / ToT 范式",
            ],
        },
        {
            "icon": "💾",
            "title": "记忆系统",
            "subtitle": "Memory System",
            "color": ACCENT_SECOND,
            "points": [
                "短期工作记忆（上下文）",
                "长期记忆（向量数据库）",
                "记忆检索与遗忘机制",
                "经验总结与知识沉淀",
            ],
        },
        {
            "icon": "🔧",
            "title": "工具调用",
            "subtitle": "Tool Use",
            "color": HIGHLIGHT,
            "points": [
                "Function Calling 机制",
                "多工具编排与组合",
                "API / 插件 / RPA 集成",
                "工具学习与自主创造",
            ],
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.2)
    gap = Inches(0.35)
    start_x = Inches(0.6)

    for i, card_data in enumerate(cards):
        x = start_x + i * (card_w + gap)
        y = Inches(1.8)

        # 卡片
        card = add_card(slide, x, y, card_w, card_h, border_color=card_data["color"])

        # 顶部色条
        add_gradient_rect(slide, x, y, card_w, Inches(0.08),
                          card_data["color"], card_data["color"], angle=0)

        # 图标
        add_text(slide, x, y + Inches(0.3), card_w, Inches(0.8),
                 card_data["icon"], size=44, color=card_data["color"],
                 align=PP_ALIGN.CENTER)

        # 标题
        add_text(slide, x, y + Inches(1.2), card_w, Inches(0.5),
                 card_data["title"], size=22, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)

        # 英文副标题
        add_text(slide, x, y + Inches(1.7), card_w, Inches(0.4),
                 card_data["subtitle"], size=11, color=TEXT_GRAY,
                 align=PP_ALIGN.CENTER, font="Arial")

        # 分割线
        add_decorative_line(slide, x + Inches(1.45), y + Inches(2.2), Inches(1),
                           color=card_data["color"], thickness=2)

        # 要点列表
        for j, point in enumerate(card_data["points"]):
            py = y + Inches(2.6 + j * 0.55)

            # 小对勾
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x + Inches(0.4), py + Inches(0.12),
                                          Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = card_data["color"]
            dot.line.fill.background()

            add_text(slide, x + Inches(0.65), py, card_w - Inches(0.8), Inches(0.4),
                     point, size=13, color=TEXT_LIGHT)

    add_footer_brand(slide)
    add_page_number(slide, 6, 12)
    return slide


# ==========================================
# 第7页：流程图页 - Agent Loop
# ==========================================
def make_flow():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "Agent 执行循环 (Agent Loop)", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "观察 → 思考 → 行动 → 再观察，经典 OODA 循环的 AI 实现", size=14, color=TEXT_GRAY)

    # 流程步骤
    steps = [
        ("接收任务", "User Task", ACCENT_PRIMARY),
        ("理解与规划", "Understand & Plan", ACCENT_PRIMARY),
        ("工具调用", "Tool Calling", ACCENT_SECOND),
        ("执行行动", "Execute Action", ACCENT_SECOND),
        ("结果观察", "Observe Result", HIGHLIGHT),
        ("反思迭代", "Reflect & Iterate", HIGHLIGHT),
    ]

    step_w = Inches(1.8)
    step_h = Inches(1.0)
    gap = Inches(0.15)
    total_w = len(steps) * step_w + (len(steps) - 1) * gap
    start_x = (SW - total_w) / 2
    y = Inches(2.8)

    for i, (name, en, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # 步骤框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h)
        shape.adjustments[0] = 0.3
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = color
        r1.font.name = "Microsoft YaHei"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = en
        r2.font.size = Pt(9)
        r2.font.color.rgb = TEXT_GRAY
        r2.font.name = "Arial"

        # 箭头（除了最后一个）
        if i < len(steps) - 1:
            arrow_x = x + step_w + Inches(0.01)
            arrow_y = y + step_h / 2 - Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            arrow_x, arrow_y, gap - Inches(0.02), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DIVIDER
            arrow.line.fill.background()

    # 下方反馈回路说明
    feedback_y = Inches(4.3)
    add_decorative_line(slide, start_x + step_w / 2, feedback_y,
                       total_w - step_w, color=ACCENT_PRIMARY, thickness=2)

    # 回环箭头文字
    add_text(slide, start_x, feedback_y + Inches(0.1), total_w, Inches(0.4),
             "↩  反馈循环：未达成目标则继续迭代，直到任务完成或触发终止条件",
             size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # 底部特性卡片
    features = [
        ("迭代次数可控", "最大步数限制，防止无限循环"),
        ("终止条件明确", "任务完成 / 失败 / 超时"),
        ("中间状态可观测", "每步行动透明可追溯"),
    ]

    feat_w = Inches(3.9)
    feat_start = Inches(0.6)
    feat_y = Inches(5.2)

    for i, (title, desc) in enumerate(features):
        x = feat_start + i * (feat_w + Inches(0.35))
        card = add_card(slide, x, feat_y, feat_w, Inches(1.4))

        # 左侧色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, feat_y, Inches(0.08), Inches(1.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [ACCENT_PRIMARY, ACCENT_SECOND, HIGHLIGHT][i]
        bar.line.fill.background()

        add_text(slide, x + Inches(0.3), feat_y + Inches(0.2), feat_w - Inches(0.4), Inches(0.5),
                 title, size=15, bold=True, color=TEXT_WHITE)
        add_text(slide, x + Inches(0.3), feat_y + Inches(0.7), feat_w - Inches(0.4), Inches(0.6),
                 desc, size=12, color=TEXT_GRAY)

    add_footer_brand(slide)
    add_page_number(slide, 7, 12)
    return slide


# ==========================================
# 第8页：数据对比页 - 效果数据
# ==========================================
def make_data_metrics():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "实践效果与关键指标", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "某企业级 Agent 平台上线 3 个月核心数据", size=14, color=TEXT_GRAY)

    # 顶部 4 个数据指标卡
    metrics = [
        ("87%", "任务完成率", "+23%", ACCENT_PRIMARY),
        ("3.2x", "效率提升", "vs 人工处理", ACCENT_SECOND),
        ("94.6%", "准确率", "↑ 5.8%", HIGHLIGHT),
        ("1200+", "日均调用量", "月增 40%", ACCENT_PRIMARY),
    ]

    m_w = Inches(2.9)
    m_h = Inches(1.8)
    m_gap = Inches(0.25)
    m_start = Inches(0.6)
    m_y = Inches(1.8)

    for i, (value, label, delta, color) in enumerate(metrics):
        x = m_start + i * (m_w + m_gap)

        card = add_card(slide, x, m_y, m_w, m_h, border_color=color)

        # 数值
        add_text(slide, x, m_y + Inches(0.2), m_w, Inches(0.9),
                 value, size=36, bold=True, color=color,
                 align=PP_ALIGN.CENTER, font="Arial")

        # 标签
        add_text(slide, x, m_y + Inches(1.05), m_w, Inches(0.4),
                 label, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

        # 增量
        add_text(slide, x, m_y + Inches(1.4), m_w, Inches(0.35),
                 delta, size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 下方 - 左右对比
    # 左：应用场景分布
    left_card = add_card(slide, Inches(0.6), Inches(3.9), Inches(6.0), Inches(2.8))
    add_text(slide, Inches(0.9), Inches(4.1), Inches(5.5), Inches(0.5),
             "▍应用场景分布", size=18, bold=True, color=ACCENT_PRIMARY)

    bars = [
        ("代码开发", 68, ACCENT_PRIMARY),
        ("数据分析", 52, ACCENT_SECOND),
        ("文档处理", 75, HIGHLIGHT),
        ("客服问答", 45, ACCENT_PRIMARY),
    ]

    for i, (name, val, color) in enumerate(bars):
        by = Inches(4.7 + i * 0.45)
        add_text(slide, Inches(0.9), by, Inches(1.5), Inches(0.4),
                 name, size=12, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)

        # 背景条
        bg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(2.3), by + Inches(0.07),
                                         Inches(4), Inches(0.28))
        bg_bar.adjustments[0] = 0.5
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(0x23, 0x2F, 0x4D)
        bg_bar.line.fill.background()

        # 进度条
        bar_w = Inches(4 * val / 100)
        fg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(2.3), by + Inches(0.07),
                                         bar_w, Inches(0.28))
        fg_bar.adjustments[0] = 0.5
        fg_bar.fill.solid()
        fg_bar.fill.fore_color.rgb = color
        fg_bar.line.fill.background()

        # 数值
        add_text(slide, Inches(5.7), by, Inches(1), Inches(0.4),
                 f"{val}%", size=11, bold=True, color=color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font="Arial")

    # 右：收益对比
    right_card = add_card(slide, Inches(6.9), Inches(3.9), Inches(5.9), Inches(2.8))
    add_text(slide, Inches(7.2), Inches(4.1), Inches(5.5), Inches(0.5),
             "▍核心收益", size=18, bold=True, color=ACCENT_SECOND)

    benefits = [
        ("人力成本", "降低 60%", "同等工作量下人力投入大幅减少"),
        ("响应速度", "提升 5 倍", "7x24 小时不间断服务"),
        ("错误率", "下降 45%", "标准化流程减少人为失误"),
    ]

    for i, (title, value, desc) in enumerate(benefits):
        by = Inches(4.7 + i * 0.65)

        # 数字
        add_text(slide, Inches(7.2), by, Inches(2), Inches(0.45),
                 value, size=18, bold=True, color=HIGHLIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)

        # 标题和描述
        add_text(slide, Inches(9.2), by, Inches(3.4), Inches(0.4),
                 title, size=13, bold=True, color=TEXT_WHITE)
        add_text(slide, Inches(9.2), by + Inches(0.35), Inches(3.4), Inches(0.3),
                 desc, size=11, color=TEXT_GRAY)

    add_footer_brand(slide)
    add_page_number(slide, 8, 12)
    return slide


# ==========================================
# 第9页：案例页 - 典型场景
# ==========================================
def make_case_study():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "典型落地场景", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "三个已验证的高价值 Agent 应用场景", size=14, color=TEXT_GRAY)

    cases = [
        {
            "tag": "研发效能",
            "title": "智能代码助手 Agent",
            "color": ACCENT_PRIMARY,
            "desc": "集成代码仓库、CI/CD、文档系统，帮助开发者完成需求分析、代码生成、Review、Bug 定位等全流程任务。",
            "highlights": ["代码评审效率 +70%", "Bug 修复时间缩短 40%", "新人上手周期减少 50%"],
        },
        {
            "tag": "数据分析",
            "title": "数据洞察 Agent",
            "color": ACCENT_SECOND,
            "desc": "自然语言驱动的数据分析助手，自动连接数据库、生成 SQL、制作图表、撰写分析报告，让数据赋能全员。",
            "highlights": ["取数需求响应 <5min", "数据分析师人效 3x", "业务自助分析率 85%"],
        },
        {
            "tag": "运营自动化",
            "title": "智能运营 Agent",
            "color": HIGHLIGHT,
            "desc": "7x24 小时自动化运营，覆盖内容生成、用户触达、活动策划、数据监控，实现千人千面的精细化运营。",
            "highlights": ["运营人力节省 60%", "用户转化率 +28%", "活动上线周期从周到小时"],
        },
    ]

    card_w = Inches(4.0)
    card_h = Inches(5.2)
    gap = Inches(0.25)
    start_x = Inches(0.55)
    y = Inches(1.8)

    for i, case in enumerate(cases):
        x = start_x + i * (card_w + gap)

        # 卡片
        card = add_card(slide, x, y, card_w, card_h, border_color=case["color"])

        # 顶部标签
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            x + Inches(0.3), y + Inches(0.3),
                                            Inches(1.4), Inches(0.4))
        tag_shape.adjustments[0] = 0.5
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = case["color"]
        tag_shape.line.fill.background()
        tf = tag_shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = case["tag"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = BG_DARK
        run.font.name = "Microsoft YaHei"

        # 标题
        add_text(slide, x + Inches(0.3), y + Inches(0.9), card_w - Inches(0.6), Inches(0.7),
                 case["title"], size=18, bold=True, color=TEXT_WHITE)

        # 分割线
        add_decorative_line(slide, x + Inches(0.3), y + Inches(1.65), Inches(0.8),
                           color=case["color"], thickness=2)

        # 描述
        add_text(slide, x + Inches(0.3), y + Inches(1.85), card_w - Inches(0.6), Inches(1.5),
                 case["desc"], size=12, color=TEXT_GRAY)

        # 效果标题
        add_text(slide, x + Inches(0.3), y + Inches(3.3), card_w - Inches(0.6), Inches(0.4),
                 "✦ 核心效果", size=13, bold=True, color=case["color"])

        # 效果列表
        for j, h in enumerate(case["highlights"]):
            hy = y + Inches(3.8 + j * 0.42)
            add_text(slide, x + Inches(0.3), hy, card_w - Inches(0.6), Inches(0.4),
                     "• " + h, size=12, color=TEXT_LIGHT)

    add_footer_brand(slide)
    add_page_number(slide, 9, 12)
    return slide


# ==========================================
# 第10页：挑战与问题页
# ==========================================
def make_challenges():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "当前挑战与应对思路", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)

    challenges = [
        {
            "title": "可靠性与幻觉",
            "icon": "⚠",
            "desc": "Agent 可能生成错误信息或执行错误操作，导致业务风险",
            "solution": "多轮验证 + 人工把关 + 沙箱执行 + 置信度评估",
            "color": HIGHLIGHT,
        },
        {
            "title": "上下文窗口限制",
            "icon": "📏",
            "desc": "长任务、复杂场景下信息丢失，影响推理连贯性",
            "solution": "分层记忆 + 摘要压缩 + 向量检索 + 滚动上下文",
            "color": ACCENT_PRIMARY,
        },
        {
            "title": "成本与延迟",
            "icon": "⏱",
            "desc": "多轮推理 + 工具调用导致响应慢、费用高",
            "solution": "大小模型协作 + 缓存机制 + 并行调用 + 本地推理",
            "color": ACCENT_SECOND,
        },
        {
            "title": "安全与可控性",
            "icon": "🔒",
            "desc": "工具调用权限、数据泄露、Prompt 注入等安全风险",
            "solution": "权限分级 + 审计日志 + 输入输出过滤 + 沙箱隔离",
            "color": HIGHLIGHT,
        },
    ]

    card_w = Inches(6.0)
    card_h = Inches(2.4)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    start_x = Inches(0.6)
    start_y = Inches(1.8)

    for i, ch in enumerate(challenges):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_card(slide, x, y, card_w, card_h, border_color=ch["color"])

        # 左侧色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ch["color"]
        bar.line.fill.background()

        # 图标 + 标题
        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.6),
                 ch["icon"], size=24, color=ch["color"], anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.9), y + Inches(0.25), card_w - Inches(1), Inches(0.55),
                 ch["title"], size=17, bold=True, color=TEXT_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

        # 问题描述
        add_text(slide, x + Inches(0.3), y + Inches(0.9), card_w - Inches(0.6), Inches(0.5),
                 "▸ 问题：" + ch["desc"], size=12, color=TEXT_GRAY)

        # 解决方案
        add_rich_text(slide, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(0.8),
                      [[
                          {"text": "▸ 思路：", "size": 12, "bold": True, "color": ch["color"]},
                          {"text": ch["solution"], "size": 12, "color": TEXT_LIGHT},
                      ]])

    add_footer_brand(slide)
    add_page_number(slide, 10, 12)
    return slide


# ==========================================
# 第11页：未来展望页
# ==========================================
def make_future():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    add_gradient_rect(slide, Inches(0), Inches(0), SW, Inches(0.08),
                      ACCENT_PRIMARY, ACCENT_SECOND)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.8),
             "未来发展方向", size=28, bold=True, color=TEXT_WHITE)
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(1),
                       color=ACCENT_PRIMARY, thickness=3)
    add_text(slide, Inches(1.8), Inches(1.2), Inches(10), Inches(0.6),
             "Agent 技术演进的五大趋势预判", size=14, color=TEXT_GRAY)

    trends = [
        ("01", "多模态 Agent", "从纯文本走向图文音视频全方位感知与交互", ACCENT_PRIMARY),
        ("02", "多 Agent 协作", "专业化分工 + 协同机制，应对复杂系统性任务", ACCENT_PRIMARY),
        ("03", "端侧 Agent", "轻量化模型 + 端侧推理，保护隐私降低延迟", ACCENT_SECOND),
        ("04", "自主进化", "持续学习 + 自我改进，Agent 能力随使用不断提升", ACCENT_SECOND),
        ("05", "Agent 经济", "Agent 之间自主交易与协作，形成新的数字经济生态", HIGHLIGHT),
    ]

    list_y = Inches(2.0)
    for i, (num, title, desc, color) in enumerate(trends):
        y = list_y + i * Inches(0.95)

        # 编号背景
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          Inches(0.8), y + Inches(0.1),
                                          Inches(0.7), Inches(0.7))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = color
        num_box.line.fill.background()
        tf = num_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = BG_DARK
        run.font.name = "Arial"

        # 标题
        add_text(slide, Inches(1.7), y + Inches(0.1), Inches(5), Inches(0.35),
                 title, size=17, bold=True, color=TEXT_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

        # 描述
        add_text(slide, Inches(1.7), y + Inches(0.45), Inches(10), Inches(0.4),
                 desc, size=13, color=TEXT_GRAY)

        # 连接线（除最后一个）
        if i < len(trends) - 1:
            line = slide.shapes.add_connector(1,
                                               Inches(1.15), y + Inches(0.8),
                                               Inches(1.15), y + Inches(0.95))
            line.line.color.rgb = DIVIDER
            line.line.width = Pt(1)

    # 右侧装饰
    add_text(slide, Inches(8.5), Inches(1.8), Inches(4.5), Inches(5),
             "🤖", size=200, color=BG_CARD, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 右侧金句
    quote_card = add_card(slide, Inches(8.2), Inches(5.5), Inches(4.5), Inches(1.2))
    add_text(slide, Inches(8.5), Inches(5.65), Inches(3.9), Inches(1),
             "\"Agent 是软件 2.0 的终极形态，\n将重塑数字世界的交互方式。\"",
             size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer_brand(slide)
    add_page_number(slide, 11, 12)
    return slide


# ==========================================
# 第12页：致谢/结束页
# ==========================================
def make_thanks():
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # 左右渐变装饰条
    add_gradient_rect(slide, Inches(0), Inches(0), Inches(0.15), SH,
                      ACCENT_PRIMARY, ACCENT_SECOND, angle=90)

    # 装饰元素 - 几何图形
    # 右上圆
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(-1),
                                 Inches(3), Inches(3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = ACCENT_PRIMARY
    c1.line.fill.background()
    sp = c1._element
    srgb = sp.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '12000')

    # 左下圆
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5.5),
                                 Inches(3), Inches(3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = ACCENT_SECOND
    c2.line.fill.background()
    sp = c2._element
    srgb = sp.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '12000')

    # 中央内容
    add_text(slide, Inches(0), Inches(2.2), SW, Inches(1.5),
             "THANK YOU", size=72, bold=True, color=TEXT_WHITE,
             align=PP_ALIGN.CENTER, font="Arial")

    # 装饰线
    add_decorative_line(slide, Inches(5.66), Inches(3.6), Inches(2),
                       color=ACCENT_PRIMARY, thickness=3)

    add_text(slide, Inches(0), Inches(3.9), SW, Inches(0.8),
             "感谢聆听 · 欢迎交流", size=24, color=TEXT_LIGHT,
             align=PP_ALIGN.CENTER)

    # 联系信息卡片
    info_card = add_card(slide, Inches(4.16), Inches(5.0), Inches(5), Inches(1.5))
    add_rich_text(slide, Inches(4.16), Inches(5.15), Inches(5), Inches(1.2),
                  [
                      [{"text": "📧  ", "size": 14, "color": ACCENT_PRIMARY},
                       {"text": "contact@example.com", "size": 13, "color": TEXT_LIGHT, "font": "Arial"}],
                      [{"text": "🌐  ", "size": 14, "color": ACCENT_PRIMARY},
                       {"text": "www.example.com", "size": 13, "color": TEXT_LIGHT, "font": "Arial"}],
                      [{"text": "💬  ", "size": 14, "color": ACCENT_PRIMARY},
                       {"text": "技术团队 · AI Agent 小组", "size": 13, "color": TEXT_LIGHT}],
                  ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_number(slide, 12, 12)
    return slide


# ==========================================
# 生成所有幻灯片
# ==========================================
make_cover()                              # 1 封面
make_toc()                                # 2 目录
make_section("01", "背景与趋势", "AI Agent 时代已然到来")  # 3 章节页
make_two_column()                         # 4 两栏内容
make_architecture()                       # 5 架构图
make_three_cards()                        # 6 三大核心技术
make_flow()                               # 7 Agent Loop
make_data_metrics()                       # 8 数据指标
make_case_study()                         # 9 落地场景
make_challenges()                         # 10 挑战
make_future()                             # 11 展望
make_thanks()                             # 12 结束页

output = "AI-Agent-技术汇报-模板.pptx"
prs.save(output)
print(f"✅ PPT 模板已生成：{output}")
print(f"📊 共 {len(prs.slides)} 页幻灯片")
