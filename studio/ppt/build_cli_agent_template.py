from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from PIL import Image
import re


OUT = "终端_AI_Agent_技术调研汇报模板.pptx"
W, H = 13.333, 7.5

# Palette: field-notebook / terminal engineering
CREAM = "F6F2E7"; PAPER = "FCFAF2"; SAND = "EEE7D6"; INK = "252B35"
MUTED = "6E6A60"; LINE = "DDD4C0"; BLUE = "2E6FA6"; ORANGE = "B85C36"
GREEN = "2E8657"; PURPLE = "765A9B"; RED = "B94A48"; WHITE = "FFFFFF"
MONO = "Avenir Next Condensed"; SANS = "PingFang SC"


def rgb(hexstr):
    hexstr = hexstr.replace('#', '')
    return RGBColor(*bytes.fromhex(hexstr))

def set_bg(slide, color=CREAM):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = rgb(color)

def rect(slide, x, y, w, h, fill=PAPER, line=None, radius=False, transparency=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill); sh.fill.transparency = transparency
    sh.line.color.rgb = rgb(line if line else fill)
    return sh

def line(slide, x1, y1, x2, y2, color=LINE, width=1, dash=None):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = rgb(color); sh.line.width = Pt(width)
    if dash: sh.line.dash_style = dash
    return sh

def text(slide, s, x, y, w, h, size=16, color=INK, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0, italic=False):
    # Presentation-first typography: avoid the small, document-like text that is
    # hard to scan from a distance. Titles retain their authored size; supporting
    # copy, tables, labels and footers receive a readable minimum.
    if size < 9:
        size += 2
    elif size <= 10.5:
        size += 4.5
    elif size <= 12:
        size += 4
    elif size <= 14:
        size += 2.5
    elif size < 20:
        size += 1.5
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = rgb(color)
    tf.vertical_anchor = valign
    return tb

def label(slide, s, x, y, w=None, color=ORANGE):
    return text(slide, s.upper(), x, y, w or 4, .22, 8.5, color, MONO, True)

def tag(slide, s, x, y, w, color=BLUE, fill=PAPER):
    rect(slide, x, y, w, .29, fill, color, True)
    text(slide, s, x, y+.035, w, .18, 7.5, color, MONO, True, PP_ALIGN.CENTER)

def picture_cover(slide, path, x, y, w, h):
    """Place an image in an editable fixed frame without aspect-ratio distortion."""
    with Image.open(path) as im:
        image_ratio = im.width / im.height
    frame_ratio = w / h
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > frame_ratio:
        crop = (1 - frame_ratio / image_ratio) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - image_ratio / frame_ratio) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic

def dot_grid(slide):
    # Intentionally empty: a clean background is easier to scan in presentation mode.
    # Information hierarchy is carried by cards, labels, and section dividers instead.
    pass

def page(slide, num, section, title, subtitle=""):
    set_bg(slide); dot_grid(slide)
    label(slide, f"{num:02d}  /  {section}", .58, .55, 3.2)
    text(slide, title, .58, .88, 11.8, .55, 26, INK, SANS, True)
    if subtitle: text(slide, subtitle, .60, 1.48, 11.2, .34, 10.5, MUTED)
    text(slide, f"TERMINAL AI AGENT RESEARCH  ·  {num:02d}", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, body, accent=BLUE, number=None):
    rect(slide, x, y, w, h, PAPER, LINE, True)
    rect(slide, x, y, .06, h, accent, accent)
    if number:
        text(slide, number, x+.23, y+.22, .5, .25, 8.5, accent, MONO, True)
        tx = x+.75
    else: tx = x+.25
    text(slide, title, tx, y+.21, w-(tx-x)-.2, .31, 13, INK, SANS, True)
    text(slide, body, x+.25, y+.63, w-.48, h-.78, 9.5, MUTED, SANS)

def add_bullets(slide, items, x, y, w, h, size=13, color=INK, accent=ORANGE):
    for i, item in enumerate(items):
        yy = y + i*.53
        rect(slide, x, yy+.11, .09, .09, accent, accent, True)
        text(slide, item, x+.19, yy, w-.19, .36, size, color)

prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# 01 Cover
s = prs.slides.add_slide(blank); set_bg(s); dot_grid(s)
rect(s, 0, 0, .16, H, ORANGE, ORANGE)
label(s, "TERMINAL AI AGENT · INTERNAL SHARING", .68, .72, 5.4)
text(s, "终端 AI Agent\n技术架构与实践", .68, 1.28, 8.1, 1.55, 34, INK, SANS, True)
text(s, "初步揭秘引发生产力巨变的 AI Agent，部分理论和实践", .72, 3.13, 8.4, .38, 15, MUTED)
rect(s, .72, 3.80, 7.55, .80, SAND, ORANGE, True)
label(s, "CHAMPION CONCLUSION", .98, 4.03, 1.85, ORANGE)
text(s, "终端 AI Agent 的效率，不来自“更快生成”，而来自可验证的「理解 → 执行 → 验证」闭环。", 2.62, 3.96, 5.30, .47, 12.2, INK, SANS, True)
text(s, "实证参考：4,867 位开发者的现场随机试验，AI 编程辅助使完成任务数 +26.08%。", .74, 4.91, 7.52, .20, 9.2, MUTED, SANS)
tag(s, "内部技术分享", .72, 5.34, 1.24, ORANGE, SAND)
tag(s, "2026.07", 2.08, 5.34, .90, BLUE, PAPER)
rect(s, 8.95, .83, 3.35, 5.55, SAND, LINE, True)
text(s, "$ terminal-ai-agent", 9.28, 1.19, 2.6, .25, 11, ORANGE, MONO, True)
for i, (a, b, c) in enumerate([("USER INTENT", "plan", BLUE), ("TOOL LOOP", "execute", ORANGE), ("VERIFICATION", "observe", GREEN), ("DELIVERY", "report", PURPLE)]):
    yy = 1.78 + i*1.02
    rect(s, 9.28, yy, 2.67, .66, PAPER, LINE, True)
    text(s, a, 9.46, yy+.14, 1.45, .15, 7.5, c, MONO, True)
    text(s, b, 9.46, yy+.34, 1.65, .18, 10.5, INK, MONO, True)
    if i < 3: text(s, "↓", 10.43, yy+.68, .25, .22, 13, MUTED, MONO, True, PP_ALIGN.CENTER)
text(s, "汇报人：刘通  /  团队：数据生态组", .72, 6.53, 4.6, .2, 9, MUTED, SANS)

# 02 Agenda
s = prs.slides.add_slide(blank); page(s, 2, "OPENING", "本次分享的三个部分", "理论概念与工程模块  →  Harness 运行时  →  团队实践与评测")
for i, (n, title, body, pages, c) in enumerate([( "01", "CLI AI Agent", "定义、核心概念、工程组件与六项 Harness 能力。", "03–07", BLUE),
                                                    ( "02", "Harness", "四个 Agent 的 Loop、Context 与 State 工程实现。", "08–12", ORANGE),
                                                    ( "03", "agent_biz", "工程实践、GAIA 对照、回放、总结与展望。", "13–21", GREEN)]):
    x = .64 + i*4.17
    rect(s, x, 2.35, 3.65, 2.75, PAPER, LINE, True)
    text(s, n, x+.26, 2.65, .52, .25, 11, c, MONO, True)
    rect(s, x+.27, 3.17, 2.77, .04, c, c)
    text(s, title, x+.27, 3.31, 2.98, .54, 15, INK, SANS, True)
    text(s, body, x+.27, 4.02, 2.96, .56, 8.0, MUTED)
    text(s, pages, x+.27, 4.84, 2.96, .20, 10, c, MONO, True)
text(s, "阅读路径：从“是什么”，到“如何运行”，再到“如何在企业内部构建、验证与持续演进”。", .66, 5.86, 11.0, .28, 10.2, MUTED, MONO)

# 03 executive
s = prs.slides.add_slide(blank); page(s, 3, "EXECUTIVE SUMMARY", "先给结论：终端 AI Agent 的价值在“可执行、可验证、可演进”", "此页用于替换成汇报中最重要的 3 条结论。")
card(s, .66, 2.18, 3.82, 2.50, "结论 01｜不是聊天界面", "核心差异是 Agent 能调用工具、观察结果，并在终端中完成多步任务。", BLUE, "01")
card(s, 4.75, 2.18, 3.82, 2.50, "结论 02｜Harness 是操作系统", "LLM 提供推理能力；Harness 管理上下文、工具、权限、恢复与验证。", ORANGE, "02")
card(s, 8.84, 2.18, 3.82, 2.50, "结论 03｜小而透明更适合演进", "以可替换的核心 Loop 起步，通过 Benchmark 驱动增强，而非堆叠框架。", GREEN, "03")
rect(s, .66, 5.28, 12.0, .72, SAND, LINE, True)
text(s, "建议决策：〔例如：以自研轻量 Harness 为核心，兼容主流工具协议，并用任务集持续校准能力〕", .95, 5.51, 11.35, .22, 11, INK, SANS, True)

# 04 definition
s = prs.slides.add_slide(blank); page(s, 4, "PRINCIPLE", "终端 AI Agent 到底是什么？", "一句话：在终端中将语言理解、工具调用和反馈循环闭合的任务执行者。")
rect(s, .66, 2.2, 5.45, 3.55, PAPER, LINE, True)
label(s, "TRADITIONAL CLI", .96, 2.52, 2.0, MUTED)
text(s, "确定性命令", .96, 2.90, 2.4, .32, 17, INK, SANS, True)
text(s, "输入参数  →  程序执行  →  固定输出", .96, 3.35, 4.35, .28, 12, MUTED, MONO)
rect(s, .94, 3.78, 4.86, 1.22, INK, INK, True)
picture_cover(s, "traditional_cli.png", .96, 3.80, 4.82, 1.18)
text(s, "适合：明确、单步、可预期的操作", .96, 5.26, 4.3, .25, 10, MUTED)
rect(s, 7.16, 2.2, 5.51, 3.55, PAPER, ORANGE, True)
label(s, "TERMINAL AI AGENT", 7.46, 2.52, 2.7, ORANGE)
text(s, "面向目标的循环", 7.46, 2.90, 3.3, .32, 17, INK, SANS, True)
text(s, "目标  →  计划  →  调用工具  →  观察  →  修正", 7.46, 3.35, 4.5, .28, 12, ORANGE, MONO, True)
rect(s, 7.44, 3.78, 4.89, 1.22, INK, INK, True)
picture_cover(s, "agent_cli.png", 7.46, 3.80, 4.85, 1.18)
text(s, "适合：开放、多步、需根据环境反馈调整的任务", 7.46, 5.26, 4.55, .25, 10, MUTED)
text(s, "关键变化：从「用户告诉程序怎么做」转向「用户说明想达成什么」。", .68, 6.18, 8.8, .27, 12, INK, SANS, True)

# 04A Engineering components
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "04A  /  ENGINEERING", .58, .55, 3.2)
text(s, "终端 AI Agent 的工程组件：从界面到真实执行", .58, .88, 11.8, .55, 25, INK, SANS, True)
text(s, "以项目实现为例：TUI 提供交互，Node.js 承载运行时，Harness 编排模型与工具，执行环境完成真实操作。", .60, 1.48, 11.8, .34, 12, MUTED)

# Main request path
rect(s, .70, 2.18, 2.28, 1.10, PAPER, BLUE, True)
label(s, "PRESENTATION", .93, 2.39, 1.2, BLUE)
text(s, "TUI · Ink / React", .93, 2.67, 1.74, .22, 12, INK, SANS, True)
text(s, "编辑输入、流式输出、进度展示", .93, 2.96, 1.75, .16, 8.8, MUTED)
text(s, "↔", 3.08, 2.60, .32, .22, 17, MUTED, MONO, True, PP_ALIGN.CENTER)

rect(s, 3.52, 2.18, 2.35, 1.10, PAPER, PURPLE, True)
label(s, "RUNTIME", 3.76, 2.39, 1.0, PURPLE)
text(s, "Node.js Runtime", 3.76, 2.67, 1.75, .22, 12, INK, SANS, True)
text(s, "启动、配置、信号、打包与分发", 3.76, 2.96, 1.75, .16, 8.8, MUTED)
text(s, "↔", 5.99, 2.60, .32, .22, 17, MUTED, MONO, True, PP_ALIGN.CENTER)

rect(s, 6.43, 2.03, 3.16, 3.56, SAND, ORANGE, True)
label(s, "CONTROL PLANE", 6.71, 2.28, 1.45, ORANGE)
text(s, "Agent Harness", 6.71, 2.57, 2.05, .26, 16, ORANGE, SANS, True)
for i, (n, title, desc) in enumerate([
    ("01", "Agent Engine", "循环、状态、退出条件"),
    ("02", "Session & Prompt", "历史、系统提示、工具定义"),
    ("03", "Compact & Verify", "压缩上下文、核验与修正"),
    ("04", "Skill Loading", "按需加载可复用操作手册"),
]):
    yy = 3.08 + i*.52
    rect(s, 6.71, yy+.06, .23, .23, ORANGE, ORANGE, True)
    text(s, n, 6.745, yy+.11, .17, .09, 6.6, WHITE, MONO, True, PP_ALIGN.CENTER)
    text(s, title, 7.08, yy, 1.43, .18, 9.7, INK, SANS, True)
    text(s, desc, 8.45, yy+.02, .83, .30, 7.8, MUTED, SANS)

text(s, "↔", 9.75, 2.60, .32, .22, 17, MUTED, MONO, True, PP_ALIGN.CENTER)
rect(s, 10.11, 2.18, 2.53, 1.10, PAPER, GREEN, True)
label(s, "REASONING", 10.38, 2.39, 1.0, GREEN)
text(s, "LLM API", 10.38, 2.67, 1.54, .22, 12, INK, SANS, True)
text(s, "理解、规划、生成工具调用", 10.38, 2.96, 1.74, .16, 8.8, MUTED)

# Tool execution branch
text(s, "Harness 将工具调用调度到真实执行环境，并把观察结果回灌模型", .74, 4.10, 5.36, .20, 10.2, MUTED, SANS, True)
line(s, 8.01, 5.60, 8.01, 6.02, ORANGE, 1.4)
line(s, 1.48, 6.02, 11.82, 6.02, ORANGE, 1.4)
tool_components = [
    ("文件系统", "read / write / list", BLUE),
    ("命令执行", "bash / child_process", ORANGE),
    ("浏览器", "Playwright", GREEN),
    ("常驻终端", "tmux session", PURPLE),
]
for i, (name, desc, c) in enumerate(tool_components):
    x = .70 + i*3.03
    line(s, x+1.08, 6.02, x+1.08, 6.18, ORANGE, 1.1)
    rect(s, x, 6.18, 2.18, .60, PAPER, LINE, True)
    text(s, name, x+.16, 6.34, .75, .18, 10, c, SANS, True)
    text(s, desc, x+.95, 6.36, 1.02, .16, 7.8, MUTED, MONO)
text(s, "TERMINAL AI AGENT RESEARCH  ·  04A", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

# 04B Core concept: LLM + Harness
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "04B  /  CORE CONCEPT", .58, .55, 3.2)
text(s, "AI Agent = LLM + Harness", .58, .88, 11.8, .55, 27, INK, SANS, True)
text(s, "模型负责生成与判断；Harness 把原始推理能力变成一个在真实环境中可控、可靠完成任务的系统。", .60, 1.48, 11.5, .34, 13, MUTED)
rect(s, .70, 2.22, 3.12, 1.30, PAPER, BLUE, True)
text(s, "LLM", 1.02, 2.53, .88, .28, 17, BLUE, MONO, True)
text(s, "推理引擎 / 原始能力", 1.02, 2.97, 2.32, .22, 12, INK, SANS, True)
text(s, "像 CPU：能执行推理，但不管理系统。", 1.02, 3.22, 2.45, .18, 9.5, MUTED)
text(s, "+", 4.12, 2.63, .48, .35, 23, MUTED, MONO, True, PP_ALIGN.CENTER)
rect(s, 4.85, 2.22, 4.28, 1.30, SAND, ORANGE, True)
text(s, "HARNESS", 5.18, 2.53, 1.7, .28, 17, ORANGE, MONO, True)
text(s, "运行时 / 控制与编排层", 5.18, 2.97, 3.05, .22, 12, INK, SANS, True)
text(s, "其中 Agent Loop 是执行内核。", 5.18, 3.22, 3.43, .18, 9.5, ORANGE, SANS, True)
text(s, "=", 9.42, 2.63, .48, .35, 23, MUTED, MONO, True, PP_ALIGN.CENTER)
rect(s, 10.13, 2.22, 2.50, 1.30, PAPER, GREEN, True)
text(s, "AI AGENT", 10.43, 2.53, 1.74, .28, 15, GREEN, MONO, True)
text(s, "可执行的任务系统", 10.43, 3.02, 1.74, .22, 11, INK, SANS, True)
text(s, "Harness 的职责：Agent Loop + 四个保障层", .72, 4.22, 5.2, .29, 15, INK, SANS, True)
os_items = [("01", "Agent Loop", "决策、调用、观察", ORANGE), ("02", "上下文与记忆", "拼装、压缩、恢复", BLUE), ("03", "工具与进程", "选择、执行、隔离", GREEN), ("04", "权限与恢复", "审批、重试、断点", RED), ("05", "验证与观测", "测试、确认、指标", PURPLE)]
for i, (n, name, desc, c) in enumerate(os_items):
    x = .72 + i*2.43
    rect(s, x, 4.73, 2.18, 1.06, PAPER, LINE, True)
    text(s, n, x+.18, 4.94, .27, .18, 9, c, MONO, True)
    text(s, name, x+.54, 4.91, 1.35, .20, 11, INK, SANS, True)
    text(s, desc, x+.18, 5.31, 1.77, .18, 9.2, MUTED)
rect(s, .72, 6.15, 11.9, .45, SAND, LINE, True)
text(s, "关系：Agent Loop ⊂ Harness —— Loop 让 Agent 不断行动；Harness 让行动可控、可靠、可恢复。", .98, 6.29, 11.25, .16, 10.5, ORANGE, SANS, True)
text(s, "TERMINAL AI AGENT RESEARCH  ·  04B", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

# 05 loop
s = prs.slides.add_slide(blank); page(s, 5, "PRINCIPLE", "Agent Loop：Harness 内的执行内核", "Harness 先提供上下文、工具与边界；Loop 再驱动 LLM 和工具反复协作，直到任务完成。")
rect(s, .68, 2.18, 12.0, 3.44, PAPER, ORANGE, True)
label(s, "AGENT LOOP / HARNESS EXECUTION KERNEL", .96, 2.42, 3.6, ORANGE)
tag(s, "Exit: task done / user stop / budget exhausted", 8.02, 2.36, 4.05, RED, PAPER)
steps = [("01", "理解目标", BLUE), ("02", "规划下一步", ORANGE), ("03", "执行工具", PURPLE), ("04", "观察结果", GREEN), ("05", "验证 / 继续", RED)]
xs = [1.0, 3.32, 5.64, 7.96, 10.28]
for i, (n, name, c) in enumerate(steps):
    rect(s, xs[i], 3.28, 1.7, .9, PAPER, c, True)
    text(s, n, xs[i]+.18, 3.47, .28, .16, 8, c, MONO, True)
    text(s, name, xs[i]+.18, 3.70, 1.26, .20, 11, INK, SANS, True)
    if i < 4:
        text(s, "→", xs[i]+1.79, 3.54, .38, .22, 16, MUTED, MONO, True, PP_ALIGN.CENTER)
line(s, 11.2, 4.56, 11.2, 5.04, ORANGE, 1.5)
line(s, 11.2, 5.04, 2.0, 5.04, ORANGE, 1.5)
line(s, 2.0, 5.04, 2.0, 4.56, ORANGE, 1.5)
text(s, "未完成：带着 observation 回到下一轮", 3.2, 5.15, 5.6, .22, 10, ORANGE, MONO, True, PP_ALIGN.CENTER)
text(s, "Harness 在每一轮提供并约束 Loop：", .82, 5.93, 2.18, .18, 9.5, MUTED, SANS, True)
tag(s, "上下文 / Prompt", 3.02, 5.86, 1.65, PURPLE, PAPER)
tag(s, "工具 / 进程", 5.06, 5.86, 1.50, GREEN, PAPER)
tag(s, "权限 / 策略", 6.95, 5.86, 1.50, RED, PAPER)
tag(s, "验证 / 恢复", 8.84, 5.86, 1.50, ORANGE, PAPER)

# 06 architecture
s = prs.slides.add_slide(blank); page(s, 6, "ARCHITECTURE", "Harness 的工程实现：1 个执行内核 + 5 个保障层", "Agent Loop 负责让任务持续向前；其余模块负责让每一步可理解、可控、可验证、可恢复。")
tag(s, "HARNESS = AGENT LOOP + SUPPORTING LAYERS", .66, 1.84, 3.65, ORANGE, PAPER)
mods = [("01", "Agent Loop", "执行内核：状态、决策迭代、退出条件", ORANGE), ("02", "Context & Prompt", "拼装任务、历史、环境与约束", BLUE), ("03", "Tool Orchestration", "定义、选择、校验与执行隔离", GREEN), ("04", "State & Memory", "压缩历史、保存状态、断点续跑", PURPLE), ("05", "Permission & Policy", "风险分级、审批、审计与边界", RED), ("06", "Verification & Recovery", "测试、确认、重试、回滚与升级", BLUE)]
for i, (n, title, body, c) in enumerate(mods):
    row, col = divmod(i, 3); x = .66 + col*4.10; y = 2.15 + row*1.78
    card(s, x, y, 3.76, 1.34, title, body, c, n)
text(s, "工程判断：先让 Loop 跑起来，再按任务风险补齐上下文、工具、状态、权限与验证；它们共同构成 Harness。", .68, 5.99, 11.3, .25, 11, INK, SANS, True)

# 06B Case studies: four concrete harness implementations
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06B  /  CASE STUDIES", .58, .55, 3.2)
text(s, "四个代表性终端 Agent：分别是谁，为什么值得看？", .58, .88, 11.8, .55, 24, INK, SANS, True)
text(s, "它们不在竞争同一个位置：从商业编码助手、个人 Agent 入口，到可复用 Harness 与自我进化的自主 Agent。", .60, 1.48, 11.8, .34, 12, MUTED)

case_studies = [
    ("CLAUDE CODE", "Anthropic", "终端里的 Agentic Coding", "商业编码 Agent 的标杆", "138k ★", BLUE,
     "让模型直接理解代码库、执行命令与 Git 工作流。"),
    ("OPENCLAW", "OpenClaw 社区", "你的个人 AI 助手：任意 OS、任意平台", "开源个人 Agent 的破圈案例", "383k ★", ORANGE,
     "复用 Pi Agent Core；自身聚焦 Gateway、渠道与编排。"),
    ("PI", "Earendil Works", "Agent Harness：可组合的运行时积木", "开发者导向的开源基础设施", "70.7k ★", GREEN,
     "把 LLM API、Agent Loop、TUI 和 Coding CLI 拆成独立组件。"),
    ("HERMES AGENT", "Nous Research", "The agent that grows with you", "自我学习 Agent 的新代表", "214k ★", PURPLE,
     "从经验中沉淀技能、跨会话记忆，并可部署到自己的环境。"),
]
for i, (name, maker, slogan, influence, stars, c, takeaway) in enumerate(case_studies):
    row, col = divmod(i, 2)
    x = .66 + col*6.12; y = 2.10 + row*2.15
    rect(s, x, y, 5.72, 1.82, PAPER, LINE, True)
    rect(s, x, y, .07, 1.82, c, c)
    text(s, name, x+.28, y+.20, 2.75, .19, 11, c, MONO, True)
    text(s, maker, x+4.22, y+.22, 1.13, .16, 8.2, MUTED, MONO, True, PP_ALIGN.RIGHT)
    text(s, slogan, x+.28, y+.52, 4.83, .23, 13, INK, SANS, True)
    text(s, influence, x+.28, y+.90, 3.42, .18, 9.2, c, SANS, True)
    text(s, stars, x+4.20, y+.85, 1.15, .25, 14, c, MONO, True, PP_ALIGN.RIGHT)
    line(s, x+.28, y+1.23, x+5.32, y+1.23, LINE, .6)
    text(s, takeaway, x+.28, y+1.42, 4.92, .23, 9.2, MUTED, SANS)
text(s, "影响力指标：GitHub Stars（截至 2026-07-14）；它衡量社区关注度，不等同于企业收入或实际部署规模。", .68, 6.55, 11.0, .20, 9, MUTED, SANS)
text(s, "TERMINAL AI AGENT RESEARCH  ·  06B", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

# 06C Harness source implementation overview
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06C  /  SOURCE MAP", .58, .55, 3.2)
text(s, "Harness 能力与源代码实现：谁在负责什么？", .58, .88, 11.8, .55, 24, INK, SANS, True)
text(s, "同一能力的责任边界不同：Claude Code 偏产品型 CLI，OpenClaw 偏 Gateway 编排，Pi 是可复用 runtime，Hermes 偏自主运行平台。", .60, 1.48, 11.8, .34, 10.5, MUTED)
source_map = [
    ("01", "Agent Loop", [("CC", "queryLoop()"), ("OC", "embedded Pi runner"), ("PI", "runLoop()"), ("HS", "run_conversation()")], ORANGE),
    ("02", "Context & Prompt", [("CC", "effective prompt + compact"), ("OC", "workspace prompt"), ("PI", "session / JSONL"), ("HS", "prompt parts + compressor")], BLUE),
    ("03", "Tool Orchestration", [("CC", "typed Tool + hooks"), ("OC", "OpenClaw tool adapter"), ("PI", "ToolDefinition"), ("HS", "registry.register()")], GREEN),
    ("04", "State & Memory", [("CC", "SessionMemory"), ("OC", "transcript + parentId"), ("PI", "session harness"), ("HS", "memory + checkpoint")], PURPLE),
    ("05", "Permission & Policy", [("CC", "tool permission check"), ("OC", "deny / allowlist / full"), ("PI", "delegate to sandbox"), ("HS", "command guards + approval")], RED),
    ("06", "Verification & Recovery", [("CC", "Stop / tool hooks"), ("OC", "result / compact guards"), ("PI", "events + extensions"), ("HS", "verify hooks + recovery")], BLUE),
]
for i, (n, title, entries, c) in enumerate(source_map):
    row, col = divmod(i, 3); x = .66 + col*4.10; y = 2.04 + row*2.05
    rect(s, x, y, 3.76, 1.65, PAPER, LINE, True); rect(s, x, y, .07, 1.65, c, c)
    text(s, n, x+.22, y+.19, .38, .22, 8.5, c, MONO, True)
    text(s, title, x+.67, y+.17, 2.78, .25, 12, INK, SANS, True)
    for j, (agent, impl) in enumerate(entries):
        yy = y+.55+j*.23
        text(s, agent, x+.23, yy, .34, .14, 7.2, c, MONO, True)
        text(s, impl, x+.63, yy-.01, 2.84, .17, 7.8, MUTED, MONO)
text(s, "阅读方式：这一页定位“源码责任归属”；后续三页用经典问题横向比较不同实现如何取舍。", .68, 6.49, 11.2, .22, 10, INK, SANS, True)
text(s, "TERMINAL AI AGENT RESEARCH  ·  06C", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

def comparison_page(code, title, subtitle, rows):
    s = prs.slides.add_slide(blank); set_bg(s)
    label(s, f"{code}  /  SOURCE COMPARISON", .58, .55, 3.8)
    text(s, title, .58, .88, 11.8, .55, 22, INK, SANS, True)
    text(s, subtitle, .60, 1.48, 11.8, .34, 10.2, MUTED)
    headers = ["经典问题", "CLAUDE CODE", "OPENCLAW", "PI", "HERMES AGENT"]
    xs = [.62, 2.45, 4.99, 7.53, 10.07]; ws = [1.75, 2.42, 2.42, 2.42, 2.42]
    for i, h in enumerate(headers):
        rect(s, xs[i], 1.96, ws[i], .42, SAND, LINE)
        text(s, h, xs[i]+.10, 2.08, ws[i]-.18, .14, 7.6, [ORANGE, BLUE, ORANGE, GREEN, PURPLE][i], MONO, True, PP_ALIGN.CENTER)
    for r, (issue, hint, cells) in enumerate(rows):
        y = 2.52 + r*1.34
        rect(s, xs[0], y, ws[0], 1.22, "F7F1E4", LINE, True)
        text(s, issue, xs[0]+.14, y+.17, ws[0]-.26, .33, 10, INK, SANS, True)
        text(s, hint, xs[0]+.14, y+.62, ws[0]-.26, .36, 7.3, MUTED, SANS)
        for c, (status, summary, source, color) in enumerate(cells):
            x = xs[c+1]
            rect(s, x, y, ws[c+1], 1.22, PAPER, LINE, True); rect(s, x, y, .045, 1.22, color, color)
            text(s, status, x+.14, y+.14, .17, .18, 8.5, color, MONO, True)
            text(s, summary, x+.36, y+.12, ws[c+1]-.52, .54, 8.1, INK, SANS, True)
            text(s, source, x+.14, y+.88, ws[c+1]-.25, .16, 6.7, MUTED, MONO)
    text(s, "✓ 明确内建   ~ 部分覆盖 / 交给上层   ✕ 未见统一的一等机制   ·   结论基于当前源码快照。", .64, 6.68, 11.6, .18, 8, MUTED, MONO)
    text(s, f"TERMINAL AI AGENT RESEARCH  ·  {code}", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)
    return s

comparison_page("06D", "四个 Agent 如何处理 Agent Loop 的三个经典问题", "共性：模型无后续工具调用则结束。差异：谁拥有停止标准、重试预算与规划责任。", [
    ("停止标准", "完成、上限与中断", [("✓", "tool-use 信号 + maxTurns、abort、Stop hook", "query.ts · queryLoop", BLUE), ("~", "Pi Loop 结束；业务完成条件交给 workspace / workflow", "pi-embedded-runner/run.ts", ORANGE), ("~", "最终文本 / 无工具调用即结束；不强制领域验收", "packages/agent/agent-loop.ts", GREEN), ("✓", "finish reason + 终止分支；verify nudge 可追加一轮", "conversation_loop.py", PURPLE)]),
    ("无限 / 无效重试", "失败分类、预算与恢复", [("✓", "maxTurns + fallback；413 压缩 / Stop hook 定向重试", "query.ts · retry transitions", BLUE), ("~", "继承 Pi；另有 compaction retry、tool-result guard", "pi-embedded-runner/*guard*", ORANGE), ("✕", "无“无进展”硬熔断；由宿主补充 task budget", "packages/agent/agent-loop.ts", RED), ("✓", "TurnRetryState + max_retries + 退避 / failover", "conversation_loop.py", PURPLE)]),
    ("Planning & Replanning", "计划表达、更新与证据", [("~", "Plan mode / Todo / Skill；模型依据结果自行重规划", "query.ts · permissionMode", ORANGE), ("✕", "无独立 planner；由 Pi prompt、Skill、workflow 承担", "embedded runner + workspace", RED), ("✕", "core 刻意无规划意见；交给宿主和外部工具", "packages/agent/agent-loop.ts", RED), ("~", "提示、记忆、Skill / 工作流承载；无强制里程碑协议", "conversation_loop.py + memory", ORANGE)])
])

comparison_page("06E", "四个 Agent 如何处理 Context & Prompt 的三个经典问题", "上下文工程决定模型能否保留目标与证据，并在 token 受限时持续推进长任务。", [
    ("超长输入", "截断、重读与 token 控制", [("✓", "micro / auto compact；413 时 strip、压缩后重试", "services/compact · query.ts", BLUE), ("✓", "Pi 压缩 + OpenClaw overflow / 结果链 guard", "pi-embedded-runner/compact*", ORANGE), ("✓", "JSONL session；自动 compaction 与会话重建", "agent-session.ts", GREEN), ("✓", "prompt parts + compressor；chunked retry", "prompt_builder.py", PURPLE)]),
    ("Long-horizon task", "目标与进度持续可见", [("~", "SessionMemory、CLAUDE.md、Skill / Todo；无强制 task state", "SessionMemory · skills", ORANGE), ("~", "workspace、渠道会话、cron、Skill；由上层编排", "agents/skills/workspace.ts", ORANGE), ("~", "parentId + JSONL 可续接；目标结构由宿主定义", "agent-session.ts", ORANGE), ("✓", "持久记忆 + checkpoint，长程恢复是一等能力", "memory/* · checkpoint*", PURPLE)]),
    ("Context Compression", "压缩保真、可恢复", [("✓", "micro / auto 两级压缩，主动与被动路径分开", "services/compact/compact.ts", BLUE), ("✓", "独立 handler + retry；保护 parentId / tool result", "pi-embedded-runner/compact*", ORANGE), ("✓", "runtime / Coding Agent 均有 transcript compaction", "agent-loop.ts · agent-session.ts", GREEN), ("✓", "压缩后重建请求；memory compressor 协作", "prompt_builder.py · loop.py", PURPLE)])
])

comparison_page("06F", "四个 Agent 如何处理 State & Memory 的三个经典问题", "长程任务的核心不是“记得更多”，而是状态可信、恢复不重放、并发不互相污染。", [
    ("崩溃恢复", "恢复后不重复执行", [("~", "会话续接为主；未见通用外部 action checkpoint 协议", "SessionMemory/*", ORANGE), ("✓", "parentId、compaction、tool-result guard 维护会话链", "pi-embedded-runner/session*", ORANGE), ("~", "JSONL 可续接；真实外部状态由宿主重新观察", "coding-agent/session*", ORANGE), ("✓", "checkpoint manager + memory provider 支持长期恢复", "checkpoint* · memory/*", PURPLE)]),
    ("陈旧 / 污染记忆", "来源、作用域与淘汰", [("~", "项目规则 / 会话记忆分层；无统一 TTL / 置信度模型", "claudemd* · SessionMemory", ORANGE), ("~", "workspace snapshot 稳定一次 run；长期策略偏上层", "agents/skills/workspace.ts", ORANGE), ("~", "JSONL / parentId 保存历史；不做事实冲突消解", "agent-session.ts", ORANGE), ("✓", "memory provider、compressor、learning graph 独立实现", "memory/* · learning_graph.py", PURPLE)]),
    ("多 Agent / 多渠道冲突", "所有权、版本与并发", [("~", "单 CLI 会话优先；无跨 Agent 状态协调层", "tools/Tool.ts · session*", ORANGE), ("✓", "Gateway 编排渠道、会话与 node host 生命周期", "gateway/* · agents/*", ORANGE), ("✕", "单 Agent runtime；multi-agent ownership 留给宿主", "packages/agent/*", RED), ("~", "支持子 Agent / 多渠道；无通用资源锁 / merge 协议", "agent/* · gateway/*", ORANGE)])
])

# 06G Our Agent: origin and code-size context
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06G  /  OUR AGENT", .58, .55, 3.2)
text(s, "agent_biz：从 Claude Code 出发，在真实业务环境中打磨的终端 Agent", .58, .88, 11.8, .52, 22, INK, SANS, True)
text(s, "主要参考 Claude Code 的产品体验与 Agent Loop；辅参考 OpenClaw、Pi、Hermes 的工程拆分；由 Claude Code vibe coding 起步，并在 Oceanus 实践中持续修改。", .60, 1.48, 11.8, .31, 9.8, MUTED)
origin = [("主要参考", "Claude Code", "终端交互、工具调用、编码 Agent 体验", BLUE), ("辅助借鉴", "OpenClaw / Pi / Hermes", "编排、可复用 runtime、记忆与恢复思路", ORANGE), ("实践打磨", "agent_biz", "在业务任务与 Oceanus 环境中持续调优", GREEN)]
for i,(tagv,title,body,c) in enumerate(origin):
    x=.66+i*4.10
    rect(s,x,2.03,3.76,1.23,PAPER,LINE,True); rect(s,x,2.03,.06,1.23,c,c)
    text(s,tagv,x+.23,2.22,1.15,.16,8,c,MONO,True)
    text(s,title,x+.23,2.51,3.15,.22,12,INK,SANS,True)
    text(s,body,x+.23,2.84,3.20,.20,8.0,MUTED,SANS)
text(s,"当前工作区代码量对比（线性尺度；TS / TSX / JS / JSX / Python；含项目测试与支撑代码，非“核心逻辑复杂度”结论）",.68,3.62,11.3,.20,9.2,MUTED,SANS,True)
locs=[("agent_biz (Qfin)",6048,GREEN),("Pi",217482,GREEN),("Claude Code",512685,BLUE),("OpenClaw",529083,ORANGE),("Hermes Agent",1617337,PURPLE)]
max_loc=1617337
for i,(name,loc,c) in enumerate(locs):
    yy=4.02+i*.47
    text(s,name,.82,yy+.11,1.60,.16,9.5,INK,SANS,True)
    rect(s,2.45,yy,6.65,.30,SAND,SAND,True)
    width=6.65*loc/max_loc
    rect(s,2.45,yy,width,.30,c,c,True)
    text(s,f"{loc:,} LOC",9.28,yy+.08,1.40,.15,9.1,c,MONO,True,PP_ALIGN.RIGHT)
    text(s,f"{loc/6048:.0f}×",10.91,yy+.08,.62,.15,8.5,MUTED,MONO,True,PP_ALIGN.RIGHT)
rect(s,.66,6.53,11.70,.32,"EDF3EF",GREEN,True)
text(s,"解读：agent_biz 用约 6K 行代码保留可运行、可改造、可审查的核心闭环；复杂能力按业务需要逐步补齐，而非一次复制成熟产品的全部体量。",.86,6.61,11.22,.15,8.4,GREEN,SANS,True,PP_ALIGN.CENTER)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06G",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

# 06H / 06I Our Agent: agent_biz (one legible screenshot per slide)
for code, title, subtitle, path, accent, takeaway in [
    ("06H", "agent_biz 执行界面：直接在 Oceanus 内部任务环境中工作", "开源地址：gitlab.daikuan.qihoo.net/biz-model/agent_biz  ·  读代码、调用工具、执行命令、汇总结果。", "ai-agent-runtime.png", GREEN,
     "核心价值：把 Agent 放到业务现场，让执行过程、工具调用和结果证据可见、可复核。"),
    ("06I", "agent_biz 个性化修改：用户可以直接调整 Agent 行为与性能", "模型、提示、工具、输出展示与行为逻辑均可通过源码和配置直接修改；调整过程本身也可被 Agent 协助完成。", "ai-agent-customization.png", ORANGE,
     "工程定位：透明、可审计的内部 Agent，降低对闭源云端工具的黑盒依赖与信息外发顾虑。"),
]:
    s = prs.slides.add_slide(blank); set_bg(s)
    label(s, f"{code}  /  OUR AGENT", .58, .55, 3.2)
    text(s, title, .58, .88, 11.8, .55, 22, INK, SANS, True)
    text(s, subtitle, .60, 1.48, 11.8, .26, 9.1, MUTED, SANS)
    # Keep the screenshot at its native wide ratio and leave a clear gutter for the takeaway.
    xx = 2.34; yy = 1.85; ww = 8.65; hh = 4.55
    rect(s, xx-.08, yy-.08, ww+.16, hh+.16, PAPER, LINE, True)
    picture_cover(s, path, xx, yy, ww, hh)
    rect(s, .66, 6.54, 11.98, .28, SAND, accent, True)
    text(s, takeaway, .82, 6.61, 11.62, .15, 8.1, INK, SANS, True, PP_ALIGN.CENTER)
    text(s, f"TERMINAL AI AGENT RESEARCH  ·  {code}", 8.65, 7.18, 3.95, .16, 7.5, MUTED, MONO, True, PP_ALIGN.RIGHT)

# 06J–06N Evaluation: GAIA benchmark results
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06J  /  EVALUATION", .58, .55, 3.2)
text(s, "GAIA Benchmark：测的不是“会不会答题”，而是 Agent 能不能完成真实任务", .58, .88, 11.8, .55, 22, INK, SANS, True)
text(s, "GAIA（General AI Assistants）由 Mialon 等在 ICLR 2024 发布：466 道人工设计题，面向推理、多模态、网页检索与工具使用。", .60, 1.48, 11.8, .30, 10.2, MUTED)
gaia_parts = [
    ("多步推理", "拆解条件、计算、交叉验证", ORANGE),
    ("外部信息", "浏览、检索、阅读网页与文档", BLUE),
    ("工具 / 多模态", "文件、图像、视频、代码和命令", GREEN),
    ("精确交付", "遵守格式、单位、答案粒度与约束", PURPLE),
]
for i, (title, body, c) in enumerate(gaia_parts):
    x=.66+(i%2)*6.08; y=2.15+(i//2)*1.55
    rect(s,x,y,5.70,1.18,PAPER,LINE,True); rect(s,x,y,.07,1.18,c,c)
    text(s,title,x+.26,y+.22,2.30,.22,13,c,SANS,True)
    text(s,body,x+.26,y+.62,4.74,.25,9.4,MUTED,SANS)
rect(s,.66,5.55,11.70,.82,SAND,LINE,True)
text(s,"为什么适合终端 AI Agent？",.92,5.78,2.55,.20,12,INK,SANS,True)
text(s,"它要求 Agent 把“模型推理”闭合成“检索 → 工具执行 → 结果核验 → 精确回答”的端到端流程。",3.36,5.78,8.35,.24,10.2,MUTED,SANS)
text(s,"来源：GAIA: A Benchmark for General AI Assistants（ICLR 2024 / arXiv:2311.12983）。本次仅选取 Level 1 + 2 的 30 题样本做本地对照。",.68,6.68,11.4,.18,8.0,MUTED,MONO)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06J",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06K  /  EVALUATION SETUP", .58, .55, 3.8)
text(s, "先把比较边界说清楚：这里既在比较模型，也在验证 Agent Harness", .58, .88, 11.8, .55, 22, INK, SANS, True)
text(s, "所有结果来自本次 GAIA 运行：同一 30 道 Level 1 + 2 题、同一 quasi-exact-match 评分口径；不是通用排行榜。", .60, 1.48, 11.8, .30, 10.2, MUTED)
rect(s,.66,2.05,5.72,3.55,PAPER,LINE,True); rect(s,.66,2.05,.08,3.55,BLUE,BLUE)
text(s,"对照组",.94,2.33,1.2,.20,12,BLUE,SANS,True)
text(s,"Claude Code（claudecode）",.94,2.72,4.7,.30,16,INK,SANS,True)
text(s,"模型：Claude Sonnet 5\n模式：直接推理 / 独立运行\n结果：19 / 30",.94,3.23,4.55,.85,10.0,MUTED,SANS)
text(s,"它是一个完整的产品型 CLI Agent，对照的是“成熟 Harness + 特定模型”的整体体验。",.94,4.66,4.58,.44,9.5,INK,SANS,True)
rect(s,6.64,2.05,5.72,3.55,PAPER,LINE,True); rect(s,6.64,2.05,.08,3.55,GREEN,GREEN)
text(s,"实验组",6.92,2.33,1.2,.20,12,GREEN,SANS,True)
text(s,"agent_biz（同一 Agent 框架）",6.92,2.72,4.7,.30,16,INK,SANS,True)
text(s,"模型：glm-5.2 / doubao-turbo-2.1 /\ndeepseek-v4-pro / deepseek-flash / doubao-mini-2.0\n方式：固定框架，只切换模型",6.92,3.23,4.74,.92,10.0,MUTED,SANS)
text(s,"因此：实验组内部的分差主要反映模型与服务稳定性；与 Claude Code 的比较反映“产品 Harness + 模型”的整体差异。",6.92,4.66,4.58,.48,9.5,INK,SANS,True)
rect(s,.66,5.92,11.70,.52,"F1E7D2",ORANGE,True)
text(s,"解释原则：不要把 agent_biz 内不同模型的结果差异误归因为框架能力；也不要把 Claude Code 的单次结果泛化为所有 Claude 模型。",.90,6.08,11.15,.16,8.9,ORANGE,SANS,True,PP_ALIGN.CENTER)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06K",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06L  /  EVALUATION RESULTS", .58, .55, 3.8)
text(s, "GAIA 样本严格口径：Claude Code 与 doubao-turbo-2.1 并列第一", .58, .88, 11.8, .55, 22, INK, SANS, True)
text(s, "严格口径：28 道可比较题；timeout 与 leak-tainted 正确答案均按错误计分。它衡量“能稳定产出可信答案”的端到端能力。", .60, 1.48, 11.8, .30, 10.0, MUTED)
scores = [("Claude Code · Sonnet 5",67.9,BLUE,"19/28"),("agent_biz · doubao-turbo-2.1",67.9,GREEN,"19/28"),("agent_biz · glm-5.2",64.3,ORANGE,"18/28"),("agent_biz · deepseek-v4-pro",60.7,GREEN,"17/28"),("agent_biz · deepseek-flash",57.1,"0E8F8F","16/28"),("agent_biz · doubao-mini-2.0",39.3,PURPLE,"11/28")]
for i,(name,score,c,frac) in enumerate(scores):
    y=2.08+i*.62
    text(s,name,.82,y+.10,3.48,.17,9.2,INK,SANS,True)
    rect(s,4.28,y,6.25,.38,SAND,SAND,True); rect(s,4.28,y,6.25*score/100,.38,c,c,True)
    text(s,f"{score:.1f}%",10.76,y+.08,.75,.18,10.2,c,MONO,True,PP_ALIGN.RIGHT)
    text(s,frac,11.58,y+.10,.55,.14,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)
rect(s,.66,5.99,11.70,.52,"EDF3EF",GREEN,True)
text(s,"关键信号：同一 agent_biz 框架切换模型后，成绩从 39.3% 到 67.9%；模型选择与服务可靠性是当前最大变量。",.90,6.15,11.18,.16,8.9,GREEN,SANS,True,PP_ALIGN.CENTER)
text(s,"数据：GAIA 本地样本；strict aligned comparison（28 题）。",.68,6.72,7.0,.15,7.5,MUTED,MONO)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06L",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06N  /  SUMMARY", .58, .55, 3.8)
text(s, "总结：从理解终端 AI Agent，到构建并验证企业内 Agent 框架", .58, .88, 11.8, .55, 21, INK, SANS, True)
text(s, "本次调研的价值不只是选型，更是明确了终端 AI Agent 的工程问题、可行实现路径与企业内落地边界。", .60, 1.48, 11.8, .30, 10.0, MUTED)
summary_items=[
    ("01", "理解了终端 AI Agent 的问题特点与开发挑战", "Agent Loop 之外，Context、Tool、State、Permission、Verification 共同决定任务能否安全、可靠、可恢复地完成。", ORANGE),
    ("02", "成功开发 agent_biz：轻代码量、兼容企业平台、逻辑透明", "主参考 Claude Code、辅参考其他 Agent 框架；以约 6K LOC 形成可在 Oceanus 运行、无黑盒暗门、可核查、可由用户自主调整迭代的框架。", GREEN),
    ("03", "配置国产模型后，取得有竞争力的 GAIA 样本结果", "在本地 GAIA 30 题样本的严格口径中，agent_biz + doubao-turbo-2.1 达到 67.9%（19/28），与 Claude Code + Sonnet 5 并列第一。", PURPLE),
]
for i,(n,title,body,c) in enumerate(summary_items):
    x=.66; y=2.08+i*1.33
    rect(s,x,y,11.70,1.10,PAPER,LINE,True); rect(s,x,y,.07,1.10,c,c)
    text(s,n,x+.24,y+.22,.48,.20,9,c,MONO,True)
    text(s,title,x+.82,y+.17,6.35,.23,12,c,SANS,True)
    text(s,body,x+.82,y+.53,9.92,.30,8.9,MUTED,SANS)
rect(s,.66,6.24,11.70,.42,"F1E7D2",ORANGE,True)
text(s,"下一阶段：让更多真实业务任务反哺 agent_biz，持续补齐安全、可靠、评测与系统集成能力。",.88,6.37,11.22,.15,8.9,ORANGE,SANS,True,PP_ALIGN.CENTER)
text(s,"口径：GAIA 为本地 30 题样本；严格对齐比较 28 题；结果用于工程验证，不代表官方总榜。",.68,6.76,8.0,.15,7.5,MUTED,MONO)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06N",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

# 06M GAIA question-level playback
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06M  /  QUESTION REPLAY", .58, .55, 3.8)
text(s, "GAIA 题目回放：Pie Menus first paper", .58, .88, 11.8, .48, 22, INK, SANS, True)
text(s, "Level 1 · 检索 + 作者消歧 + 文献时间线。它直观展示：同一个 Harness 换模型，及不同 Agent 产品，对“精确证据链”的处理差异。", .60, 1.42, 11.8, .27, 9.8, MUTED)
rect(s,.66,1.88,11.70,.86,"F1E7D2",LINE,True)
text(s,"题目（译）：在 2015 年论文《Pie Menus or Linear Menus, Which Is Better?》的作者中，找出此前发表过论文的那位作者；其最早发表论文的标题是什么？",.92,2.12,11.12,.32,10.2,INK,SANS,True)
rect(s,.66,2.93,11.70,.48,"EDF3EF",GREEN,True)
text(s,"标准答案：Mapping Human Oriented Information to Software Agents for Online Systems Usage",.90,3.08,11.16,.16,8.8,GREEN,MONO,True,PP_ALIGN.CENTER)
answers = [
    ("Claude Code · Sonnet 5", "✕", "Effectiveness of Mapping Human-Oriented Information to Software Agents for Online Teaching Environments", RED),
    ("agent_biz · glm-5.2", "✓", "Mapping human-oriented information to software agents for online systems usage", GREEN),
    ("agent_biz · doubao-turbo-2.1", "✓", "Mapping human-oriented information to software agents for online systems usage", GREEN),
    ("agent_biz · deepseek-v4-pro", "✕", "A new software agent 'learning' algorithm", RED),
    ("agent_biz · deepseek-flash", "✕", "Effectiveness of Mapping Human-Oriented Information to Software Agents for Online Teaching Environments", RED),
    ("agent_biz · doubao-mini-2.0", "✕", "Effectiveness of Mapping Human-Oriented Information to Software Agents for Online Teaching Environments", RED),
]
for i,(name,status,answer,c) in enumerate(answers):
    row,col=divmod(i,2); x=.66+col*6.08; y=3.67+row*.79
    rect(s,x,y,5.70,.65,PAPER,LINE,True); rect(s,x,y,.055,.65,c,c)
    text(s,status,x+.18,y+.18,.22,.18,11,c,MONO,True)
    text(s,name,x+.49,y+.13,2.45,.15,8.2,INK,SANS,True)
    text(s,answer,x+.49,y+.35,4.90,.17,6.8,MUTED,MONO)
rect(s,.66,6.18,11.70,.44,SAND,LINE,True)
text(s,"读法：这是“精确文献追溯”而非泛泛相关性判断；6 个运行中仅 glm-5.2 与 doubao-turbo-2.1 命中精确标题。",.88,6.31,11.18,.15,8.5,INK,SANS,True,PP_ALIGN.CENTER)
text(s,"数据来源：GAIA · Pie Menus first paper · Level 1。",.68,6.75,7.0,.15,7.4,MUTED,MONO)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06M",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

# 07 routes
s = prs.slides.add_slide(blank); page(s, 7, "LANDSCAPE", "技术路线对比：框架能力与工程取舍", "表格中的项目可替换为你实际调研对象；用“适配场景”代替泛泛评价。")
headers = ["路线 / 项目", "核心范式", "优势", "代价 / 风险", "适配场景"]
rows = [("轻量自研", "最小 Loop + 工具协议", "透明、可控、易裁剪", "需自行补齐工程能力", "产品早期 / 强定制"),
        ("CLI 成熟方案", "终端原生 Agent", "交互、权限、上下文完善", "抽象较深 / 依赖绑定", "研发效能场景"),
        ("通用 Agent 框架", "编排 / 多 Agent", "生态丰富、集成快", "复杂度与调试成本高", "复杂工作流验证"),
        ("云端托管 Agent", "平台化运行时", "运维与可观测性较完整", "数据、成本、可移植性", "快速试点 / 标准流程")]
colx=[.66, 2.83, 5.0, 7.4, 9.9]; widths=[2.12,2.08,2.32,2.42,2.72]
for i,h in enumerate(headers):
    rect(s,colx[i],2.15,widths[i],.52,SAND,LINE)
    text(s,h,colx[i]+.12,2.32,widths[i]-.2,.15,8.5,ORANGE if i==0 else MUTED,MONO,True)
for r,row in enumerate(rows):
    yy=2.67+r*.77
    for i,val in enumerate(row):
        rect(s,colx[i],yy,widths[i],.77,PAPER,LINE)
        text(s,val,colx[i]+.12,yy+.14,widths[i]-.2,.45,9.2,INK if i==0 else MUTED,SANS,i==0)
text(s,"调研口径：统一选取〔任务类型〕、〔模型〕、〔环境〕，避免把模型能力误归因于框架。",.68,6.15,10.8,.24,10,MUTED,MONO)

# 08 benchmark
s = prs.slides.add_slide(blank); page(s, 8, "EVALUATION", "Benchmark 不是分数表，而是能力拆解工具", "推荐把任务按“是否需要真实执行 / 是否可验证 / 风险等级”三个维度切分。")
items = [("A", "单步工具调用", "读取文件、搜索代码、格式转换", BLUE), ("B", "多步任务规划", "定位问题 → 修改 → 执行测试", ORANGE), ("C", "环境反馈修正", "命令失败后分析原因并调整", GREEN), ("D", "安全与权限", "危险命令确认、敏感信息处理", RED)]
for i,(a,t,b,c) in enumerate(items):
    x=.66+(i%2)*6.1; y=2.16+(i//2)*1.6
    rect(s,x,y,5.72,1.2,PAPER,LINE,True)
    text(s,a,x+.25,y+.27,.35,.22,13,c,MONO,True)
    text(s,t,x+.82,y+.21,3.2,.25,13,INK,SANS,True)
    text(s,b,x+.82,y+.60,4.45,.25,9.5,MUTED)
text(s,"评分建议：任务成功率 40%  ·  过程可靠性 30%  ·  安全合规 20%  ·  成本与耗时 10%",.68,5.74,10.8,.26,11,INK,MONO,True)

# 09 results
s = prs.slides.add_slide(blank); page(s, 9, "EVALUATION", "测试结果模板：让结论一眼可读", "仅展示统一条件下的结果，并在页脚补充任务数、模型版本、运行环境。")
methods=[("方案 A",BLUE,72),("方案 B",ORANGE,84),("我们的 Agent",GREEN,91)]
for i,(name,c,score) in enumerate(methods):
    yy=2.28+i*1.12
    text(s,name,.78,yy+.13,1.45,.22,11,INK,SANS,True)
    rect(s,2.38,yy,7.35,.48,SAND,SAND,True)
    rect(s,2.38,yy,7.35*score/100,.48,c,c,True)
    text(s,str(score)+"%",9.94,yy+.11,.75,.22,12,c,MONO,True)
text(s,"可替换图表：成功率 / 平均耗时 / 成本 / 人工介入次数 / 安全拒绝率",.78,5.86,8.5,.25,10,MUTED)
rect(s,9.0,2.13,3.36,2.86,PAPER,LINE,True)
label(s,"READOUT",9.28,2.43,1.0,ORANGE)
text(s,"结果解读",9.28,2.79,2.1,.27,14,INK,SANS,True)
add_bullets(s,["领先项：〔填写〕","短板项：〔填写〕","下一轮验证：〔填写〕"],9.28,3.30,2.55,1.2,9.5,INK,ORANGE)
text(s,"N=〔任务数量〕  ·  Model=〔版本〕  ·  Env=〔环境〕",.78,6.57,6.6,.18,8.5,MUTED,MONO)

# 10 product
s = prs.slides.add_slide(blank); page(s, 10, "OUR AGENT", "我们的 Agent：一个可读、可改、可执行的核心", "此页用于介绍项目定位和最能支撑“为什么做”的 3 个特性。")
card(s,.66,2.20,3.75,2.58,"逻辑透明", "核心状态机与工具调用路径可追踪；新同学可以从主循环快速理解系统。", BLUE,"01")
card(s,4.80,2.20,3.75,2.58,"核心简洁", "保持最小依赖与少量抽象，复杂能力按需以插件或策略层叠加。", ORANGE,"02")
card(s,8.94,2.20,3.75,2.58,"执行能力强", "面向真实仓库与终端环境，强调观察、纠错、验证的闭环完成率。", GREEN,"03")
text(s,"项目入口：〔repo / command / 文档链接〕",.68,5.42,6.3,.24,10,ORANGE,MONO,True)
text(s,"适用边界：〔目标用户〕  ·  〔典型任务〕  ·  〔当前不处理的场景〕",.68,5.90,10.6,.23,10,MUTED)

# 11 code scale
s = prs.slides.add_slide(blank); page(s, 11, "OUR AGENT", "复杂度对比：用代码量解释，别用代码量证明", "LOC 只作为理解维护成本的辅助证据，需搭配能力覆盖和结果数据。")
bars=[("我们的核心 Loop", "〔120〕 LOC", 2.6, GREEN), ("方案 A 关键路径", "〔680〕 LOC", 6.3, BLUE), ("方案 B 关键路径", "〔1,240〕 LOC", 9.4, ORANGE)]
for i,(name,v,l,c) in enumerate(bars):
    yy=2.32+i*1.15
    text(s,name,.82,yy+.13,2.0,.22,10.5,INK,SANS,True)
    rect(s,3.12,yy,l,.48,SAND,SAND,True)
    rect(s,3.12,yy,l*.78,.48,c,c,True)
    text(s,v,3.12+l+.16,yy+.12,1.0,.2,10,c,MONO,True)
rect(s,.82,5.91,11.54,.53,PAPER,LINE,True)
text(s,"口径提示：统计范围 = 〔核心 Agent 路径〕；排除 〔生成代码 / 依赖 / 测试〕；日期 = 〔YYYY-MM-DD〕。",1.04,6.08,10.9,.18,9,MUTED,MONO)

# 12 roadmap
s = prs.slides.add_slide(blank); page(s, 12, "NEXT", "下一步：用测试集驱动能力迭代", "路线图按“先闭环、再加固、后规模化”组织。")
phases=[("NOW", "闭环可用", ["稳定核心 Loop", "沉淀 20+ 代表任务", "定义成功与退出标准"], BLUE),
        ("NEXT", "工程加固", ["权限策略与审计", "上下文压缩与恢复", "验证器 / 测试沙箱"], ORANGE),
        ("LATER", "规模化演进", ["插件和工具生态", "可观测性与评估平台", "团队协同与知识沉淀"], GREEN)]
for i,(tagv,title,bs,c) in enumerate(phases):
    x=.66+i*4.12
    rect(s,x,2.23,3.72,3.38,PAPER,LINE,True)
    tag(s,tagv,x+.24,2.50,.74,c,SAND if c==ORANGE else PAPER)
    text(s,title,x+.25,3.03,2.8,.29,15,INK,SANS,True)
    add_bullets(s,bs,x+.26,3.63,3.05,1.4,9.5,INK,c)
    if i<2: text(s,"→",x+3.78,3.72,.3,.2,16,MUTED,MONO,True,PP_ALIGN.CENTER)

# 13 closing
s = prs.slides.add_slide(blank); set_bg(s); dot_grid(s)
label(s,"DISCUSSION",.68,.80,1.5,ORANGE)
text(s,"把终端 AI Agent 当作\n一条可验证的工程闭环。",.68,1.38,8.5,1.28,29,INK,SANS,True)
text(s,"建议从一个真实任务集开始：先让它能完成，再让它可靠、可控、可持续演进。",.72,3.28,8.8,.32,13,MUTED)
rect(s,.72,4.33,7.65,.80,SAND,LINE,True)
text(s,"讨论题：我们的第一个“必须稳定完成”的任务是什么？",1.02,4.60,6.95,.22,12,ORANGE,SANS,True)
text(s,"THANK YOU",.72,6.45,2.0,.18,9,MUTED,MONO,True)
text(s,"〔联系人 / 项目地址〕",9.03,6.45,3.55,.18,9,MUTED,MONO,True,PP_ALIGN.RIGHT)

# Keep the evaluation narrative in this order: setup → score → question replay → readout.
def find_slide_index(needle):
    for i, slide in enumerate(prs.slides):
        if any(hasattr(shape, "text") and needle in shape.text for shape in slide.shapes):
            return i
    raise ValueError(f"Slide not found: {needle}")

score_idx = find_slide_index("GAIA 样本严格口径")
replay_idx = find_slide_index("Pie Menus first paper")
if replay_idx != score_idx + 1:
    sld_id_list = prs.slides._sldIdLst
    replay_id = sld_id_list[replay_idx]
    sld_id_list.remove(replay_id)
    sld_id_list.insert(score_idx + 1, replay_id)

# Remove the generic template tail after the GAIA readout, then add the final two pages.
readout_idx = find_slide_index("总结：从理解终端 AI Agent")

# 06O Outlook
s = prs.slides.add_slide(blank); set_bg(s)
label(s, "06O  /  OUTLOOK", .58, .55, 3.2)
text(s, "后续展望：让终端 AI Agent 成为企业业务模型迭代的可靠基础设施", .58, .88, 11.8, .52, 22, INK, SANS, True)
text(s, "方向不是追求一次性“全自动”，而是把 Agent Loop、业务系统和评测治理逐步连成可验证、可回退的闭环。", .60, 1.48, 11.8, .28, 10.0, MUTED)
outlook = [
    ("01", "继续打磨企业内 Agent Loop + Auto Research", "补齐安全、可靠、权限、验证、恢复与审计；把自动研究能力用于各类业务模型迭代。", ORANGE),
    ("02", "关联更多业务系统，提升自动化", "从信息读取走向问题归因、问题抽象、动作建议、执行验证；优先接入可控、可回滚的系统。", GREEN),
    ("03", "建立人机协同的评测与治理飞轮", "沉淀真实任务集、成功/可靠/成本/风险指标和证据链；将人工审批与复盘变成下一轮能力输入。", PURPLE),
]
for i,(n,title,body,c) in enumerate(outlook):
    x=.66+i*4.10
    rect(s,x,2.20,3.76,2.20,PAPER,LINE,True); rect(s,x,2.20,.07,2.20,c,c)
    text(s,n,x+.23,2.48,.45,.20,9,c,MONO,True)
    text(s,title,x+.23,2.85,3.18,.50,12,INK,SANS,True)
    text(s,body,x+.23,3.53,3.18,.48,8.8,MUTED,SANS)
rect(s,1.34,5.12,10.66,.92,"EDF3EF",GREEN,True)
text(s,"实践中的失败、人工介入、权限拒绝、超时与验收证据",1.55,5.34,3.25,.20,9.2,INK,SANS,True,PP_ALIGN.CENTER)
text(s,"→",5.00,5.34,.38,.20,16,GREEN,MONO,True,PP_ALIGN.CENTER)
text(s,"反哺 agent_biz：工具、提示、状态、权限与验证机制持续迭代",5.34,5.34,4.10,.20,9.2,GREEN,SANS,True,PP_ALIGN.CENTER)
text(s,"→",9.61,5.34,.38,.20,16,GREEN,MONO,True,PP_ALIGN.CENTER)
text(s,"更适配企业业务",10.02,5.34,1.70,.20,9.2,INK,SANS,True,PP_ALIGN.CENTER)
text(s,"原则：每扩展一项自动化能力，都同步定义它的责任边界、评测任务、人工接管点和回滚方案。",.68,6.48,11.3,.20,9.2,INK,SANS,True,PP_ALIGN.CENTER)
text(s,"TERMINAL AI AGENT RESEARCH  ·  06O",8.65,7.18,3.95,.16,7.5,MUTED,MONO,True,PP_ALIGN.RIGHT)

# 06P Closing
s = prs.slides.add_slide(blank); set_bg(s); dot_grid(s)
rect(s,0,0,.16,H,ORANGE,ORANGE)
label(s,"THANK YOU / Q&A",.68,.80,2.3,ORANGE)
text(s,"从一个真实任务开始，\n把 Agent 做成可验证的企业能力。",.68,1.40,9.40,1.18,28,INK,SANS,True)
text(s,"欢迎讨论：第一个值得交给终端 AI Agent 稳定完成的业务任务，应该是什么？",.72,3.25,9.70,.30,13,MUTED,SANS)
rect(s,.72,4.25,8.40,.82,SAND,LINE,True)
text(s,"感谢聆听  ·  Questions & Discussion",1.04,4.56,7.72,.22,13,ORANGE,SANS,True,PP_ALIGN.CENTER)
text(s,"agent_biz / Qfin",.72,6.45,3.2,.18,9,MUTED,MONO,True)
text(s,"〔联系人 / 项目地址〕",9.03,6.45,3.55,.18,9,MUTED,MONO,True,PP_ALIGN.RIGHT)

# Keep the two new closing pages, while dropping the older generic template pages.
sld_id_list = prs.slides._sldIdLst
while len(sld_id_list) > readout_idx + 3:
    sld_id_list.remove(sld_id_list[readout_idx + 1])

# Drop the former executive-summary template page; it is superseded by the agenda and final summary.
executive_idx = find_slide_index("先给结论：终端 AI Agent")
prs.slides._sldIdLst.remove(prs.slides._sldIdLst[executive_idx])

# Renumber visible section labels and presentation footers after every insertion/removal.
for slide_no, slide in enumerate(prs.slides, 1):
    page_no = f"{slide_no:02d}"
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = re.sub(r"^\d{2}[A-Z]?(?=\s*/)", page_no, run.text)
                run.text = re.sub(r"(?<=TERMINAL AI AGENT RESEARCH  ·  )\d{2}[A-Z]?\b", page_no, run.text)

prs.save(OUT)
print(OUT)
